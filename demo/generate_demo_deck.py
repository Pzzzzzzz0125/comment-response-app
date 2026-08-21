#!/usr/bin/env python3
"""Generate a nontechnical coworker demo deck for Permit Precedents."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("Permit_Precedents_Coworker_Demo.pptx")

W = Inches(13.333)
H = Inches(7.5)

TEAL = RGBColor(0x06, 0x5F, 0x59)
TEAL_DARK = RGBColor(0x0B, 0x3B, 0x39)
TEAL_LIGHT = RGBColor(0xE7, 0xF5, 0xF2)
MINT = RGBColor(0xD5, 0xF2, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x24, 0x26)
MUTED = RGBColor(0x5E, 0x6B, 0x6D)
BORDER = RGBColor(0xD8, 0xDE, 0xDC)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_LIGHT = RGBColor(0xFF, 0xF6, 0xDD)
GREEN = RGBColor(0x15, 0x80, 0x3D)
GREEN_LIGHT = RGBColor(0xE9, 0xF8, 0xED)
BG = RGBColor(0xF7, 0xF8, 0xF6)


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_text(slide, text, x, y, w, h, size=22, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_text(slide, parts, x, y, w, h, size=20, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    p = frame.paragraphs[0]
    p.alignment = align
    for text, bold, part_color in parts:
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = part_color or color
    return box


def add_card(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(1)
    return card


def add_pill(slide, text, x, y, w, fill=TEAL_LIGHT, color=TEAL, border=BORDER):
    pill = add_card(slide, x, y, w, 0.42, fill, border)
    add_text(slide, text, x + 0.12, y + 0.08, w - 0.24, 0.25, 11, color, True,
             align=PP_ALIGN.CENTER)
    return pill


def add_header(slide, kicker, title, subtitle=None, number=None):
    add_text(slide, kicker.upper(), 0.65, 0.42, 6.2, 0.28, 12, TEAL, True)
    add_text(slide, title, 0.65, 0.83, 11.9, 0.72, 30, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.55, 11.8, 0.55, 15, MUTED)
    if number is not None:
        add_text(slide, str(number), 12.15, 0.42, 0.45, 0.28, 11, MUTED, True,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, text="Permit Precedents · Coworker Demo"):
    add_text(slide, text, 0.65, 7.16, 6.0, 0.18, 9, MUTED)


def add_circle_icon(slide, label, x, y, color=TEAL, fill=TEAL_LIGHT, size=0.62):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.color.rgb = color
    circle.line.width = Pt(1.5)
    add_text(slide, label, x, y + 0.12, size, 0.32, 14, color, True, align=PP_ALIGN.CENTER)


def add_bullet_list(slide, bullets, x, y, w, line_h=0.62, size=18):
    for i, item in enumerate(bullets):
        yy = y + i * line_h
        add_circle_icon(slide, "✓", x, yy + 0.02, size=0.36, fill=MINT)
        add_text(slide, item, x + 0.55, yy, w - 0.55, line_h, size, INK)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL_DARK)
    add_pill(slide, "COMMENT RESPONSE APP", 0.72, 0.62, 2.15, MINT, TEAL_DARK, MINT)
    add_text(slide, "Permit\nPrecedents", 0.72, 1.42, 7.0, 1.9, 48, WHITE, True)
    add_text(slide, "Find prior review history. Understand what happened. Verify the original evidence.",
             0.76, 3.5, 6.7, 1.05, 22, RGBColor(0xD6, 0xEA, 0xE6))
    add_text(slide, "Coworker demo", 0.76, 6.52, 3.0, 0.3, 14, WHITE, True)

    panel = add_card(slide, 8.05, 0.78, 4.55, 5.95, WHITE, WHITE)
    add_text(slide, "A faster path from a question to evidence", 8.48, 1.22, 3.7, 0.75, 22, INK, True)
    items = [
        ("1", "Choose a city", "Keep the search relevant to the jurisdiction."),
        ("2", "Ask a question", "Use normal work language."),
        ("3", "Review the history", "See comments, responses, and recurring issues."),
        ("4", "Open the source", "Verify the exact page, paragraph, or cell."),
    ]
    for i, (n, title, desc) in enumerate(items):
        yy = 2.15 + i * 1.03
        add_circle_icon(slide, n, 8.47, yy, size=0.52)
        add_text(slide, title, 9.18, yy - 0.01, 2.85, 0.3, 15, INK, True)
        add_text(slide, desc, 9.18, yy + 0.31, 2.82, 0.42, 11, MUTED)
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Why it helps", "Permit history is valuable—but difficult to reuse", number=2)
    add_footer(slide)
    cards = [
        ("Scattered files", "Comments and responses live across PDFs, spreadsheets, Word files, projects, and review rounds.", "F"),
        ("Repeated review", "The same design issue may appear again with a new clarification or follow-up.", "R"),
        ("Slow verification", "Even after finding a useful precedent, someone still needs the exact source and location.", "V"),
    ]
    for i, (title, body, icon) in enumerate(cards):
        x = 0.65 + i * 4.22
        add_card(slide, x, 2.22, 3.75, 3.65)
        add_circle_icon(slide, icon, x + 0.32, 2.56, size=0.64)
        add_text(slide, title, x + 0.32, 3.42, 3.05, 0.44, 20, INK, True)
        add_text(slide, body, x + 0.32, 4.08, 3.05, 1.25, 15, MUTED)
    add_text(slide, "Permit Precedents turns that history into a searchable, verifiable workspace.",
             0.65, 6.25, 12.0, 0.45, 21, TEAL_DARK, True, align=PP_ALIGN.CENTER)
    return slide


def slide_capabilities(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "What users can do", "One workspace for overview, questions, records, and evidence", number=3)
    add_footer(slide)
    data = [
        ("City overview", "See the amount and shape of the history available for a selected city.", TEAL_LIGHT, TEAL),
        ("Knowledge chat", "Ask what happened before and receive an evidence-linked historical answer.", GREEN_LIGHT, GREEN),
        ("Historical Library", "Search and filter comments, responses, projects, rounds, and review status.", AMBER_LIGHT, AMBER),
        ("Issue history", "Follow one design concern through comments, responses, and reviewer follow-ups.", RGBColor(0xF0,0xEC,0xFA), RGBColor(0x6B,0x4E,0x9B)),
    ]
    for i, (title, body, fill, color) in enumerate(data):
        row, col = divmod(i, 2)
        x = 0.65 + col * 6.15
        y = 2.18 + row * 2.18
        add_card(slide, x, y, 5.68, 1.72, fill, color)
        add_text(slide, f"0{i+1}", x + 0.28, y + 0.28, 0.6, 0.35, 12, color, True)
        add_text(slide, title, x + 1.0, y + 0.26, 4.1, 0.38, 19, INK, True)
        add_text(slide, body, x + 1.0, y + 0.78, 4.18, 0.66, 14, MUTED)
    return slide


def slide_chat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Ask Permit History", "Start broad, then move to the exact supporting record", number=4)
    add_footer(slide)
    add_card(slide, 0.65, 2.08, 12.05, 4.55)
    add_pill(slide, "USER QUESTION", 0.98, 2.42, 1.55)
    add_text(slide, "How have we handled tree-related comments?", 2.72, 2.42, 6.95, 0.38, 19, INK, True)

    add_text(slide, "Historical pattern", 0.98, 3.18, 2.25, 0.3, 13, TEAL, True)
    add_text(slide,
             "Across the available projects, tree comments were generally handled through concrete plan revisions and tree-removal documentation, rather than narrative responses alone.",
             0.98, 3.58, 10.75, 0.88, 17, INK)
    add_rich_text(slide, [
        ("At 4155 Mitzi Dr, the plans and Tree Summary Data were updated to identify trees and proposed removals. ", False, INK),
        ("[1][2]", True, TEAL),
        (" At 701 S Clover Ave, tree circumferences were added to Sheet A1.01. ", False, INK),
        ("[3]", True, TEAL),
    ], 0.98, 4.63, 10.75, 1.0, 16)
    add_text(slide, "4 relevant comments · 2 projects · 3 review rounds · 4 confirmed responses",
             0.98, 5.9, 9.9, 0.28, 12, MUTED, True)
    add_pill(slide, "Click a citation to inspect the evidence", 9.65, 5.72, 2.55, TEAL_LIGHT, TEAL)
    return slide


def slide_library(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Historical Library", "Browse issues directly and compare the comment with the recorded response", number=5)
    add_footer(slide)
    add_card(slide, 0.65, 2.05, 4.05, 4.72, WHITE)
    add_text(slide, "ISSUES", 0.94, 2.34, 1.2, 0.25, 12, TEAL, True)
    add_text(slide, "Search and filters", 0.94, 2.76, 2.5, 0.34, 18, INK, True)
    for i, (title, status, fill, color) in enumerate([
        ("Tree labels and removals", "Has response", GREEN_LIGHT, GREEN),
        ("Rated wall assembly", "Review history", TEAL_LIGHT, TEAL),
        ("Drainage plan", "No response", AMBER_LIGHT, AMBER),
    ]):
        yy = 3.35 + i * 0.98
        add_card(slide, 0.94, yy, 3.45, 0.76, BG)
        add_text(slide, title, 1.12, yy + 0.13, 2.05, 0.24, 12, INK, True)
        add_pill(slide, status, 3.15, yy + 0.16, 1.0, fill, color)

    add_card(slide, 4.98, 2.05, 7.72, 2.12, TEAL_LIGHT, TEAL)
    add_text(slide, "GOVERNMENT COMMENT", 5.33, 2.4, 2.5, 0.25, 11, TEAL, True)
    add_text(slide, "Identify tree labels, circumferences, and proposed removals on the site plan.",
             5.33, 2.88, 6.7, 0.72, 18, INK, True)
    add_card(slide, 4.98, 4.38, 7.72, 2.39, GREEN_LIGHT, GREEN)
    add_text(slide, "COMPANY RESPONSE", 5.33, 4.72, 2.5, 0.25, 11, GREEN, True)
    add_text(slide, "The Existing and Demolition Site Plan and Tree Summary Data were updated.",
             5.33, 5.18, 6.55, 0.72, 18, INK, True)
    add_pill(slide, "Confirmed response", 10.45, 6.0, 1.72, GREEN_LIGHT, GREEN)
    return slide


def slide_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Recurring issue history", "See how one specific design issue changed over time", number=6)
    add_footer(slide)
    x_line = 1.35
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_line), Inches(2.18), Inches(0.025), Inches(3.9))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
    events = [
        ("Government comment", "PC1", "Reviewer states the initial requirement.", TEAL, TEAL_LIGHT),
        ("Applicant response", "PC1", "The project team identifies the revision or submitted material.", GREEN, GREEN_LIGHT),
        ("Reviewer follow-up", "PC2", "The reviewer clarifies what remains unresolved.", AMBER, AMBER_LIGHT),
        ("Later response", "PC3", "A more specific correction is recorded and linked to its source.", TEAL, TEAL_LIGHT),
    ]
    for i, (label, rnd, body, color, fill) in enumerate(events):
        yy = 2.2 + i * 1.02
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.18), Inches(yy + 0.25), Inches(0.36), Inches(0.36))
        circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.color.rgb = WHITE
        add_card(slide, 1.8, yy, 10.2, 0.88, fill, color)
        add_pill(slide, label, 2.08, yy + 0.2, 1.6, fill, color)
        add_pill(slide, rnd, 3.85, yy + 0.2, 0.62, WHITE, color)
        add_text(slide, body, 4.75, yy + 0.19, 6.75, 0.42, 14, INK)
    add_text(slide, "Repeated appearances in different files stay attached as sources—not duplicate events.",
             1.8, 6.44, 10.2, 0.36, 16, TEAL_DARK, True, align=PP_ALIGN.CENTER)
    return slide


def slide_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Verify before you rely", "Every useful answer should lead back to the original evidence", number=7)
    add_footer(slide)
    add_bullet_list(slide, [
        "Open the supporting comment and response beside the AI answer.",
        "See all source files attached to the same historical event.",
        "Open PDFs, spreadsheets, and Word previews inside the app.",
        "Navigate to the cited page, paragraph, sheet, cell, or range.",
        "Use the explicit Download Original action only when the file itself is needed.",
    ], 0.75, 2.16, 6.0, line_h=0.72, size=17)

    add_card(slide, 7.18, 2.05, 5.5, 4.78, WHITE)
    add_text(slide, "SOURCE EVIDENCE", 7.56, 2.4, 2.2, 0.28, 12, TEAL, True)
    add_card(slide, 7.56, 2.91, 4.73, 0.73, TEAL_LIGHT, TEAL)
    add_text(slide, "Primary source: Review Comments.xlsx", 7.8, 3.13, 4.2, 0.26, 14, INK, True)
    add_card(slide, 7.56, 3.83, 4.73, 0.73, BG, BORDER)
    add_text(slide, "Also appears in: Response Letter.pdf", 7.8, 4.05, 4.2, 0.26, 14, INK, True)
    add_text(slide, "Cited location", 7.56, 4.92, 1.5, 0.25, 12, MUTED, True)
    add_text(slide, "Review Comments · cell C6", 7.56, 5.3, 3.4, 0.28, 16, INK)
    add_pill(slide, "Open original source", 7.56, 5.9, 2.25, TEAL, WHITE, TEAL)
    return slide


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL_DARK)
    add_text(slide, "The goal", 0.75, 0.72, 2.2, 0.36, 14, MINT, True)
    add_text(slide, "Spend less time searching.\nSpend more time understanding.",
             0.75, 1.38, 8.7, 1.72, 38, WHITE, True)
    add_text(slide,
             "Choose a city → Ask a question → Review the history → Verify the source",
             0.78, 3.55, 10.8, 0.52, 21, RGBColor(0xD6, 0xEA, 0xE6))
    add_card(slide, 0.75, 4.55, 11.85, 1.3, RGBColor(0x10, 0x4B, 0x47), RGBColor(0x2B, 0x6B, 0x66))
    add_text(slide, "What would you want to find first?", 1.1, 4.92, 5.1, 0.42, 23, WHITE, True)
    add_text(slide, "A prior response · A recurring issue · A project precedent · A source document",
             6.05, 4.94, 5.95, 0.42, 15, MINT, align=PP_ALIGN.RIGHT)
    add_text(slide, "Live app: http://localhost:8010/", 0.78, 6.65, 4.5, 0.3, 13, WHITE, True)
    add_text(slide, "Questions & feedback", 9.4, 6.65, 3.15, 0.3, 13, WHITE, True, align=PP_ALIGN.RIGHT)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "Permit Precedents — Coworker Demo"
    prs.core_properties.subject = "User-focused demonstration of the Comment Response App"
    prs.core_properties.author = "Permit Precedents Team"
    prs.core_properties.keywords = "permit history, comments, responses, evidence, demo"

    slide_title(prs)
    slide_problem(prs)
    slide_capabilities(prs)
    slide_chat(prs)
    slide_library(prs)
    slide_timeline(prs)
    slide_sources(prs)
    slide_close(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
